"""
카카오톡 선물하기 데이터 분석 PPT 생성 스크립트
Output: kakao_gift_analysis_selfmade.pptx (21 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── 상수 ───────────────────────────────────────────────
SLIDE_W = 12192000   # 16:9 width (EMU)
SLIDE_H = 6858000    # 16:9 height (EMU)

# 카카오 브랜드 컬러
KAKAO_YELLOW = RGBColor(0xFF, 0xE8, 0x12)
BLACK        = RGBColor(0x00, 0x00, 0x00)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK    = RGBColor(0x33, 0x33, 0x33)
GRAY_MID     = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_BLUE  = RGBColor(0x1A, 0x56, 0xDB)
ACCENT_GREEN = RGBColor(0x03, 0x7A, 0x48)
ACCENT_RED   = RGBColor(0xC8, 0x1E, 0x1E)

FONT_NAME = "맑은 고딕"

BASE_DIR   = "C:/Users/user/Desktop/pjt/portfolio/kakao_gift"
CHART_DIR  = f"{BASE_DIR}/selfmade/analysis/charts"
OUTPUT     = f"{BASE_DIR}/kakao_gift_analysis_selfmade.pptx"

# ─── 유틸리티 함수 ────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    return prs

def blank_slide(prs):
    """완전히 빈 레이아웃 슬라이드 추가"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    """슬라이드 배경색 채우기"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_multiline_text(slide, lines, left, top, width, height,
                        font_size=16, bold=False, color=BLACK,
                        align=PP_ALIGN.LEFT, line_spacing=1.2):
    """여러 줄 텍스트 (리스트로 전달)"""
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, is_bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = clr
    return txBox

def add_image(slide, img_path, left, top, width, height):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Emu(left), Emu(top), Emu(width), Emu(height))
        return True
    return False

def add_header_bar(slide, part_label, slide_title, part_color=KAKAO_YELLOW):
    """상단 헤더 바: 좌측 Part 레이블 + 제목"""
    BAR_H = 700000  # ~0.76인치
    # 좌측 컬러 바
    add_rect(slide, 0, 0, 1_200_000, BAR_H, fill_color=part_color)
    add_text(slide, part_label,
             60000, 120000, 1_100_000, 500000,
             font_size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    # 제목 영역
    add_rect(slide, 1_200_000, 0, SLIDE_W - 1_200_000, BAR_H, fill_color=BLACK)
    add_text(slide, slide_title,
             1_300_000, 130000, SLIDE_W - 1_400_000, 500000,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def add_3col_frame(slide, q_text, data_lines, action_lines, top=800000):
    """
    비즈니스 질문 | 데이터 결과 | 실행 제언  3단 구조
    각 col: (left, width)
    """
    CONTENT_H = SLIDE_H - top - 200000
    COL_W = (SLIDE_W - 200000) // 3
    PAD = 80000

    cols = [
        (100000,  "💬 비즈니스 질문", GRAY_LIGHT,  ACCENT_BLUE,  q_text),
        (100000 + COL_W + 50000, "📊 데이터 결과", GRAY_LIGHT, ACCENT_GREEN, data_lines),
        (100000 + (COL_W + 50000)*2, "✅ 실행 제언", GRAY_LIGHT, ACCENT_RED, action_lines),
    ]

    for left, header, bg, hdr_color, body in cols:
        # 배경 박스
        add_rect(slide, left, top, COL_W, CONTENT_H,
                 fill_color=bg, line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=Pt(0.5))
        # 헤더
        add_rect(slide, left, top, COL_W, 380000, fill_color=hdr_color)
        add_text(slide, header,
                 left + PAD, top + 100000, COL_W - PAD*2, 280000,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        # 본문
        if isinstance(body, str):
            add_text(slide, body,
                     left + PAD, top + 420000, COL_W - PAD*2, CONTENT_H - 450000,
                     font_size=14, color=GRAY_DARK, align=PP_ALIGN.LEFT)
        elif isinstance(body, list):
            add_multiline_text(slide, body,
                               left + PAD, top + 420000, COL_W - PAD*2, CONTENT_H - 450000,
                               align=PP_ALIGN.LEFT)


# ─── 슬라이드 생성 함수 ──────────────────────────────────

def slide_01_cover(prs):
    """슬라이드 1: 표지"""
    sl = blank_slide(prs)
    fill_bg(sl, BLACK)

    # 카카오 옐로우 왼쪽 사이드 바
    add_rect(sl, 0, 0, 500000, SLIDE_H, fill_color=KAKAO_YELLOW)

    # 메인 타이틀
    add_text(sl, "카카오톡 선물하기",
             700000, 1_400_000, SLIDE_W - 800000, 900000,
             font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(sl, "데이터 분석 프로젝트",
             700000, 2_250_000, SLIDE_W - 800000, 700000,
             font_size=36, bold=False, color=KAKAO_YELLOW, align=PP_ALIGN.LEFT)

    add_text(sl, "성장 지속을 위한 사용자 행동 기반 전략 제언",
             700000, 3_050_000, SLIDE_W - 800000, 500000,
             font_size=20, bold=False, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)

    # 구분선
    add_rect(sl, 700000, 3_700_000, 2_000_000, 60000, fill_color=KAKAO_YELLOW)

    # 메타 정보
    add_text(sl, "분석 기간: 2023–2024   |   5 Phase End-to-End 분석   |   Selfmade Version",
             700000, 3_900_000, SLIDE_W - 800000, 400000,
             font_size=15, bold=False, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.LEFT)


def slide_02_problem(prs):
    """슬라이드 2: 문제 정의"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "INTRO", "왜 이 분석을 했는가? — 문제 정의")

    # 서비스 성장 현황 박스
    add_rect(sl, 100000, 820000, SLIDE_W - 200000, 700000,
             fill_color=RGBColor(0xFF, 0xF8, 0xD0), line_color=KAKAO_YELLOW, line_width=Pt(1.5))
    add_text(sl, "📈  카카오톡 선물하기 현황",
             250000, 860000, 3_000_000, 350000,
             font_size=16, bold=True, color=BLACK)
    add_text(sl, "2024 GMV ₩75.2억  (+29.7% YoY)   |   국내 모바일 선물 시장 1위   |   수신→발신 전환율 97.9%",
             250000, 1_130_000, SLIDE_W - 500000, 320000,
             font_size=15, bold=False, color=GRAY_DARK)

    # 3가지 핵심 질문
    add_text(sl, "BUT — 성장 이면의 3가지 핵심 질문",
             100000, 1_680_000, SLIDE_W - 200000, 400000,
             font_size=18, bold=True, color=BLACK)

    questions = [
        ("①  시즌 의존도",
         "GMV 성장이 이벤트 시즌에만 집중된 건 아닌가?\n비시즌 기간 사용자는 어디에 있는가?",
         "→  Phase 1 EDA",  ACCENT_BLUE),
        ("②  고객 구조",
         "어떤 사용자가 실제 비즈니스를 떠받치는가?\nCRM 예산을 어디에 써야 ROI가 높은가?",
         "→  Phase 2·3",   ACCENT_GREEN),
        ("③  바이럴 루프",
         "수신자가 발신자로 전환되는 구조가\n얼마나 효율적이고, 어디서 끊기는가?",
         "→  Phase 4·5",   ACCENT_RED),
    ]

    BOX_W = (SLIDE_W - 300000) // 3
    for i, (title, body, phase_tag, color) in enumerate(questions):
        lft = 100000 + i * (BOX_W + 50000)
        add_rect(sl, lft, 2_150_000, BOX_W, 2_150_000 + 2_000_000 - 2_150_000,
                 fill_color=GRAY_LIGHT, line_color=color, line_width=Pt(2))
        # 컬러 탑 바
        add_rect(sl, lft, 2_150_000, BOX_W, 320000, fill_color=color)
        add_text(sl, title,
                 lft + 80000, 2_200_000, BOX_W - 160000, 260000,
                 font_size=15, bold=True, color=WHITE)
        add_text(sl, body,
                 lft + 80000, 2_520_000, BOX_W - 160000, 700000,
                 font_size=14, color=GRAY_DARK)
        add_text(sl, phase_tag,
                 lft + 80000, 3_300_000, BOX_W - 160000, 280000,
                 font_size=13, bold=True, color=color)

    add_text(sl, "이 세 질문에 end-to-end로 답하는 것이 이 분석의 목적입니다.",
             100000, SLIDE_H - 450000, SLIDE_W - 200000, 350000,
             font_size=14, bold=False, color=GRAY_MID, align=PP_ALIGN.CENTER)


def slide_03_design(prs):
    """슬라이드 3: 분석 설계"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "INTRO", "분석 설계 — 5 Phase End-to-End")

    phases = [
        ("Phase 1", "EDA &\n시즌 분석",    "시즌 의존도\n진단",           KAKAO_YELLOW, BLACK),
        ("Phase 2", "RFM\n세그멘테이션",   "고가치 고객\n식별",           RGBColor(0xFF,0xA0,0x00), BLACK),
        ("Phase 3", "LTV\n코호트",         "고객 장기\n가치 산출",        RGBColor(0x03,0x7A,0x48), WHITE),
        ("Phase 4", "바이럴\n루프",        "K-factor &\n전환 메커니즘",   RGBColor(0x1A,0x56,0xDB), WHITE),
        ("Phase 5", "CRM\n전략",           "A/B·ROAS·\n액션플랜",        RGBColor(0x77,0x17,0xAA), WHITE),
    ]

    BOX_W = 2_000_000
    BOX_H = 1_600_000
    TOP   = 1_000_000
    GAP   = 100000
    TOTAL_W = BOX_W * 5 + GAP * 4
    START = (SLIDE_W - TOTAL_W) // 2

    for i, (phase, title, biz_q, bg, fg) in enumerate(phases):
        lft = START + i * (BOX_W + GAP)
        add_rect(sl, lft, TOP, BOX_W, BOX_H, fill_color=bg)
        add_text(sl, phase,
                 lft + 80000, TOP + 120000, BOX_W - 160000, 300000,
                 font_size=14, bold=True, color=fg, align=PP_ALIGN.CENTER)
        add_text(sl, title,
                 lft + 80000, TOP + 420000, BOX_W - 160000, 450000,
                 font_size=18, bold=True, color=fg, align=PP_ALIGN.CENTER)
        add_text(sl, biz_q,
                 lft + 80000, TOP + 950000, BOX_W - 160000, 500000,
                 font_size=13, color=fg, align=PP_ALIGN.CENTER)
        # 화살표 (마지막 제외)
        if i < 4:
            arr_lft = lft + BOX_W + 10000
            add_rect(sl, arr_lft, TOP + BOX_H//2 - 60000, GAP - 20000, 120000,
                     fill_color=RGBColor(0xCC, 0xCC, 0xCC))

    # 데이터 소스 표
    add_text(sl, "데이터 소스",
             100000, TOP + BOX_H + 300000, 2_000_000, 350000,
             font_size=16, bold=True, color=BLACK)

    tables = [
        ("orders", "주문 트랜잭션"),
        ("users", "사용자 정보"),
        ("gifts", "선물 발송/수신"),
        ("campaigns", "캠페인 메타"),
        ("events", "사용자 이벤트"),
    ]
    T_TOP = TOP + BOX_H + 700000
    T_W   = (SLIDE_W - 200000) // 5
    for i, (tbl, desc) in enumerate(tables):
        lft = 100000 + i * T_W
        add_rect(sl, lft, T_TOP, T_W - 40000, 650000,
                 fill_color=GRAY_LIGHT, line_color=RGBColor(0xDD,0xDD,0xDD), line_width=Pt(0.5))
        add_text(sl, f"`{tbl}`",
                 lft + 60000, T_TOP + 80000, T_W - 160000, 280000,
                 font_size=14, bold=True, color=ACCENT_BLUE)
        add_text(sl, desc,
                 lft + 60000, T_TOP + 340000, T_W - 160000, 250000,
                 font_size=12, color=GRAY_DARK)


def slide_04_gmv(prs):
    """슬라이드 4: GMV 성장 구조"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 1 · Phase 1", "성장 구조 — GMV 추이 & 트렌드 분해")
    add_3col_frame(sl,
        q_text="성장이 이벤트 시즌에만\n의존하는 건 아닌가?\n\n비시즌 기간에도\n실질적 성장이 있는가?",
        data_lines=[
            ("2023 GMV  ₩58.0억", 16, True, BLACK),
            ("2024 GMV  ₩75.2억  (+29.7%)", 16, True, ACCENT_GREEN),
            ("", 8, False, BLACK),
            ("STL 분해 결과:", 14, True, GRAY_DARK),
            ("• Trend 성분: 독립적 우상향", 14, False, GRAY_DARK),
            ("• Seasonal: 3개 피크 명확", 14, False, GRAY_DARK),
            ("• Residual: 특이값 없음", 14, False, GRAY_DARK),
        ],
        action_lines=[
            ("비시즌 투자 정당화", 15, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("시즌 성과가 좋아도", 13, False, GRAY_DARK),
            ("Trend 성분이 독립 성장 중", 13, False, GRAY_DARK),
            ("→ 비시즌 CRM 예산 확보 근거", 13, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("연간 GMV 예측 모델 구축 가능", 13, False, GRAY_DARK),
        ]
    )


def slide_05_season(prs):
    """슬라이드 5: 시즌 패턴"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 1 · Phase 1", "시즌 패턴 — 시즌별 전략을 달리해야 하는가?")
    add_3col_frame(sl,
        q_text="시즌별 구매 패턴이\n다른가?\n\n동일한 캠페인 전략으로\n모든 시즌을 커버할 수 있나?",
        data_lines=[
            ("시즌별 AOV (평균 주문금액):", 14, True, GRAY_DARK),
            ("", 6, False, BLACK),
            ("빼빼로데이  ₩40,946", 15, True, ACCENT_BLUE),
            ("→ 소액다량: 저단가 대규모 발송", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("설날  ₩128,266", 15, True, ACCENT_GREEN),
            ("→ 고단가: 프리미엄 선물세트", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("발렌타인  뷰티 76.7% 집중", 14, False, GRAY_DARK),
        ],
        action_lines=[
            ("시즌별 전략 분리 필수", 15, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("빼빼로데이:", 13, True, ACCENT_BLUE),
            ("발송량 최대화, 저단가 번들", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("설날:", 13, True, ACCENT_GREEN),
            ("프리미엄 큐레이션, 고가 묶음", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("발렌타인:", 13, True, RGBColor(0xAA,0x00,0x88)),
            ("뷰티 카테고리 집중 프로모션", 13, False, GRAY_DARK),
        ]
    )


def slide_06_category(prs):
    """슬라이드 6: 카테고리 믹스"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 1 · Phase 1", "카테고리 믹스 — 특정 카테고리 의존 리스크")
    add_3col_frame(sl,
        q_text="특정 카테고리에\n지나치게 의존하는 건 아닌가?\n\n시즌이 바뀌면\n매출 구조가 무너지는가?",
        data_lines=[
            ("카테고리 × 시즌 교차 분석:", 14, True, GRAY_DARK),
            ("", 6, False, BLACK),
            ("뷰티/화장품: 연평균 ~35%", 14, False, GRAY_DARK),
            ("  → 발렌타인 76.7% 극단 집중", 13, False, ACCENT_RED),
            ("", 6, False, BLACK),
            ("식품/음료: 빼빼로데이 강세", 14, False, GRAY_DARK),
            ("생활용품: 설날 시즌 강세", 14, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("시즌마다 지배 카테고리 교체", 14, True, ACCENT_BLUE),
        ],
        action_lines=[
            ("카테고리 다각화 로드맵", 15, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("뷰티 의존도 리스크 관리:", 13, True, GRAY_DARK),
            ("비시즌 뷰티 수요 유지 캠페인", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("크로스카테고리 번들 기획:", 13, True, GRAY_DARK),
            ("식품+뷰티 콜라보 선물세트", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("시즌별 카테고리 큐레이션 자동화", 13, False, GRAY_DARK),
        ]
    )


def slide_07_rfm(prs):
    """슬라이드 7: RFM 세그멘테이션"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 2 · Phase 2", "RFM 세그멘테이션 — 어떤 사용자가 비즈니스를 만드는가?")

    # 설계 근거 박스
    add_rect(sl, 100000, 820000, SLIDE_W//2 - 150000, 600000,
             fill_color=RGBColor(0xFF,0xF8,0xD0), line_color=KAKAO_YELLOW, line_width=Pt(1))
    add_text(sl, "🔧 도메인 특화 설계: 표준 RFM과의 차이",
             200000, 860000, SLIDE_W//2 - 300000, 280000,
             font_size=14, bold=True, color=BLACK)
    add_text(sl, "선물하기는 연간 구매 빈도가 낮음 → F 스코어를 5분위 대신 4레벨로 커스텀 재설계",
             200000, 1_120_000, SLIDE_W//2 - 300000, 280000,
             font_size=13, color=GRAY_DARK)

    # 9개 세그먼트 그리드
    segments = [
        # (이름, R범위, F범위, 색상, 설명)
        ("Champions",        "R≥4", "F≥3", RGBColor(0xFF,0xE8,0x12), BLACK),
        ("Loyal Customers",  "R=3", "F≥3", RGBColor(0xFF,0xA0,0x00), BLACK),
        ("At Risk",          "R≤2", "F≥3", ACCENT_RED,               WHITE),
        ("Potential Loyalists","R=4","F=2", ACCENT_GREEN,             WHITE),
        ("Need Attention",   "R=3", "F=2", RGBColor(0x77,0x17,0xAA), WHITE),
        ("Can't Lose Them",  "R≤2", "F≥4", RGBColor(0xC8,0x1E,0x1E), WHITE),
        ("Recent Customers", "R=4", "F=1", ACCENT_BLUE,              WHITE),
        ("Promising",        "R=3", "F=1", RGBColor(0x02,0x62,0x8A), WHITE),
        ("About to Sleep",   "R≤2", "F=1", GRAY_MID,                 WHITE),
    ]

    COL_N = 3
    BOX_W = (SLIDE_W - 200000) // COL_N - 40000
    BOX_H = 680000
    TOP   = 1_520_000

    for i, (name, r, f, bg, fg) in enumerate(segments):
        col = i % COL_N
        row = i // COL_N
        lft = 100000 + col * (BOX_W + 60000)
        top = TOP + row * (BOX_H + 40000)
        add_rect(sl, lft, top, BOX_W, BOX_H, fill_color=bg)
        add_text(sl, name,
                 lft + 60000, top + 80000, BOX_W - 120000, 280000,
                 font_size=14, bold=True, color=fg)
        add_text(sl, f"{r}  ·  {f}",
                 lft + 60000, top + 360000, BOX_W - 120000, 240000,
                 font_size=12, color=fg)


def slide_08_pareto(prs):
    """슬라이드 8: GMV Pareto"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 2 · Phase 2", "GMV Pareto — CRM 예산을 어디에 써야 하는가?")
    add_3col_frame(sl,
        q_text="모든 세그먼트를\n동등하게 관리해야 하나?\n\nCRM 예산 배분의\n근거는 무엇인가?",
        data_lines=[
            ("세그먼트별 GMV 기여:", 14, True, GRAY_DARK),
            ("", 6, False, BLACK),
            ("Champions      44.3%", 15, True, BLACK),
            ("Loyal Customers  18.7%", 14, False, GRAY_DARK),
            ("At Risk          12.1%", 14, False, GRAY_DARK),
            ("나머지 6개       24.9%", 14, False, GRAY_MID),
            ("", 8, False, BLACK),
            ("상위 2개 세그먼트 = GMV 63%", 15, True, ACCENT_RED),
            ("Pareto 법칙 확인", 13, False, GRAY_DARK),
        ],
        action_lines=[
            ("집중 투자 전략", 15, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("Champions:", 13, True, BLACK),
            ("VIP 혜택, 얼리 액세스", 13, False, GRAY_DARK),
            ("유지 비용 최소화", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("At Risk / Can't Lose:", 13, True, ACCENT_RED),
            ("긴급 재활성화 예산 집중", 13, False, GRAY_DARK),
            ("개인화 메시지 + 쿠폰", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("Recent: 온보딩 시퀀스", 13, False, GRAY_DARK),
        ]
    )


def slide_09_ltv(prs):
    """슬라이드 9: LTV 커브"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 2 · Phase 3", "LTV 커브 — 신규 고객 획득에 얼마까지 써도 되는가?")
    add_3col_frame(sl,
        q_text="신규 고객 1명 획득에\n얼마까지 투자해도 되나?\n\nCAC (고객 획득 비용) 상한은\n데이터로 어떻게 정할 수 있나?",
        data_lines=[
            ("코호트 LTV (M+N 누적):", 14, True, GRAY_DARK),
            ("", 6, False, BLACK),
            ("M+0   ₩75,248", 14, False, GRAY_DARK),
            ("M+3   ₩112,400", 14, False, GRAY_DARK),
            ("M+6   ₩155,700", 14, False, GRAY_DARK),
            ("M+11  ₩203,169  ★", 16, True, ACCENT_GREEN),
            ("", 8, False, BLACK),
            ("M+0 대비 2.7배 성장", 14, True, ACCENT_GREEN),
        ],
        action_lines=[
            ("CAC 상한선 설정", 15, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("LTV:CAC = 10:1 목표 기준:", 13, True, GRAY_DARK),
            ("CAC ≤ ₩20,317 허용", 14, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("세그먼트별 차등 적용:", 13, True, GRAY_DARK),
            ("Champions 예상 LTV 높음", 13, False, GRAY_DARK),
            ("→ 더 높은 CAC 허용 가능", 13, False, GRAY_DARK),
            ("", 6, False, BLACK),
            ("마케팅 예산 승인 근거 확보", 13, False, GRAY_DARK),
        ]
    )


def slide_10_cohort(prs):
    """슬라이드 10: 시즌 코호트 품질"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 2 · Phase 3", "시즌 코호트 품질 — 시즌 마케팅 ROI 정당화")
    add_3col_frame(sl,
        q_text="시즌 캠페인으로 유입된\n고객은 금방 떠나는 건 아닌가?\n\n시즌 마케팅 예산 증가를\n데이터로 정당화할 수 있나?",
        data_lines=[
            ("시즌 vs 비시즌 코호트 비교:", 14, True, GRAY_DARK),
            ("", 6, False, BLACK),
            ("시즌 유입  M+11 LTV  ₩201,500", 14, False, GRAY_DARK),
            ("비시즌 유입 M+11 LTV  ₩204,800", 14, False, GRAY_DARK),
            ("", 8, False, BLACK),
            ("차이: -1.6%  →  사실상 동일", 15, True, ACCENT_GREEN),
            ("", 8, False, BLACK),
            ("M+1 리텐션도 시즌=비시즌", 13, False, GRAY_DARK),
        ],
        action_lines=[
            ("시즌 예산 확대 정당화", 15, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("시즌 유입 ≠ 일회성 고객", 14, True, ACCENT_GREEN),
            ("장기 고객 확보 채널로 입증", 13, False, GRAY_DARK),
            ("", 8, False, BLACK),
            ("실무 적용:", 13, True, GRAY_DARK),
            ("'시즌 캠페인은 비용이다'", 13, False, GRAY_DARK),
            ("→ '시즌 캠페인은 투자다'로", 13, False, GRAY_DARK),
            ("프레임 전환 가능", 13, False, GRAY_DARK),
        ]
    )


def slide_11_kfactor(prs):
    """슬라이드 11: K-factor"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 3 · Phase 4", "K-factor — 선물하기의 자체 바이럴 성장 엔진")
    add_3col_frame(sl,
        q_text="선물하기 플랫폼이\n자체적으로 사용자를\n늘리는 구조인가?\n\nK-factor는 얼마이며\n어떤 유저가 더 높은가?",
        data_lines=[
            ("K-factor (수정값):", 14, True, GRAY_DARK),
            ("", 6, False, BLACK),
            ("전체 평균    K = 1.559  ★", 16, True, ACCENT_GREEN),
            ("", 6, False, BLACK),
            ("자기선물 유저  K = 2.090", 14, True, ACCENT_BLUE),
            ("타인선물 유저  K = 1.515", 14, False, GRAY_DARK),
            ("", 8, False, BLACK),
            ("K > 1  →  바이럴 성장 가능 구조", 14, True, ACCENT_GREEN),
            ("", 6, False, BLACK),
            ("* DESA 원본 K=3.948은 중복", 12, False, GRAY_MID),
            ("  집계 오류 → 수정 적용", 12, False, GRAY_MID),
        ],
        action_lines=[
            ("자기선물 유저 집중 육성", 15, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("자기선물 K=2.090:", 13, True, ACCENT_BLUE),
            ("헤비유저 특성, 바이럴 증폭 효과", 13, False, GRAY_DARK),
            ("→ 자기선물 전용 혜택 강화", 13, False, GRAY_DARK),
            ("", 8, False, BLACK),
            ("K-factor 시뮬레이션:", 13, True, GRAY_DARK),
            ("1,000명 활성화 → M+6 ~5,000명", 13, False, GRAY_DARK),
            ("CAC 절감 효과 정량화", 13, False, GRAY_DARK),
        ]
    )


def slide_12_reciprocity(prs):
    """슬라이드 12: 상호성 & Golden Time"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 3 · Phase 4", "상호성 & Golden Time — 수신자를 언제 발신자로 전환시키는가?")
    add_3col_frame(sl,
        q_text="선물 받은 사람이\n다시 선물을 보내게 되는 건\n언제, 얼마나 일어나는가?\n\n리마인더 캠페인 타이밍은?",
        data_lines=[
            ("Pay-it-forward 비율:", 14, True, GRAY_DARK),
            ("수신자 중 발신 경험 있음  99.4%", 15, True, ACCENT_GREEN),
            ("", 8, False, BLACK),
            ("Golden Time (수신→발신 중앙값):", 14, True, GRAY_DARK),
            ("30일  ★", 18, True, ACCENT_BLUE),
            ("", 6, False, BLACK),
            ("수신 후 30일 이내 전환율 최고", 13, False, GRAY_DARK),
            ("30일 이후 전환율 급격히 감소", 13, False, GRAY_DARK),
        ],
        action_lines=[
            ("D+7/14/30 자동 넛지", 15, True, ACCENT_RED),
            ("", 8, False, BLACK),
            ("수신 후 자동 캠페인 시퀀스:", 13, True, GRAY_DARK),
            ("D+7:  '받은 선물 어때요?'", 13, False, GRAY_DARK),
            ("D+14: '당신도 선물해보세요'", 13, False, GRAY_DARK),
            ("D+30: '마지막 기회' 긴급 넛지", 13, False, GRAY_DARK),
            ("", 8, False, BLACK),
            ("Golden Time 내 전환 집중으로", 13, False, GRAY_DARK),
            ("CAC 없는 신규 발신자 확보", 13, False, GRAY_DARK),
        ]
    )


def slide_13_occasion(prs):
    """슬라이드 13: Occasion 트리거 (신규)"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 3 · Phase 4", "Occasion 트리거 — 전환은 시간이 아닌 '계기'로 발동된다")

    # 핵심 인사이트 강조 박스
    add_rect(sl, 100000, 820000, SLIDE_W - 200000, 550000,
             fill_color=RGBColor(0xFF,0xF0,0xF0), line_color=ACCENT_RED, line_width=Pt(2))
    add_text(sl, "💡  핵심 발견: 미전환의 63.2%는 '시간 부족'이 아닌 '명확한 계기(Occasion) 부재'가 원인",
             250000, 920000, SLIDE_W - 500000, 380000,
             font_size=15, bold=True, color=ACCENT_RED)

    # Occasion 순위 (좌)
    LEFT_W = SLIDE_W // 2 - 150000
    add_text(sl, "Occasion별 전환율 순위",
             100000, 1_500_000, LEFT_W, 300000,
             font_size=16, bold=True, color=BLACK)

    occasions = [
        ("🥇  생일 (Birthday)", "가장 강력한 트리거", KAKAO_YELLOW, BLACK),
        ("🥈  빼빼로데이",       "시즌 이벤트 집중", RGBColor(0xFF,0xA0,0x00), BLACK),
        ("🥉  크리스마스",       "연말 시즌",       ACCENT_GREEN, WHITE),
        ("4위  발렌타인데이",     "뷰티 카테고리",   ACCENT_BLUE, WHITE),
    ]
    for i, (name, desc, bg, fg) in enumerate(occasions):
        top = 1_870_000 + i * 700000
        add_rect(sl, 100000, top, LEFT_W, 580000, fill_color=bg)
        add_text(sl, name, 180000, top + 80000, LEFT_W - 200000, 280000,
                 font_size=15, bold=True, color=fg)
        add_text(sl, desc, 180000, top + 330000, LEFT_W - 200000, 200000,
                 font_size=13, color=fg)

    # 실행 제언 (우)
    RIGHT_L = SLIDE_W // 2 + 100000
    RIGHT_W = SLIDE_W // 2 - 200000
    add_text(sl, "✅  실행 제언 — Occasion-Triggered CRM",
             RIGHT_L, 1_500_000, RIGHT_W, 300000,
             font_size=16, bold=True, color=BLACK)

    actions = [
        ("생일 D-30 캠페인",
         "'○○님 생일에 자신에게 선물해보세요'\n→ 자기선물 K=2.090 극대화"),
        ("기념일 전 맞춤 큐레이션",
         "수신 이력 기반 개인화 상품 추천\n→ '지난번 받은 것과 비슷한 선물'"),
        ("63.2% 미전환 타겟",
         "Occasion 없는 사용자에게\n가상 계기 생성 (시즌 이벤트 연결)"),
    ]
    for i, (title, body) in enumerate(actions):
        top = 1_870_000 + i * 950000
        add_rect(sl, RIGHT_L, top, RIGHT_W, 820000,
                 fill_color=GRAY_LIGHT, line_color=ACCENT_RED, line_width=Pt(1))
        add_text(sl, title, RIGHT_L + 80000, top + 80000, RIGHT_W - 160000, 280000,
                 font_size=14, bold=True, color=ACCENT_RED)
        add_text(sl, body, RIGHT_L + 80000, top + 360000, RIGHT_W - 160000, 380000,
                 font_size=13, color=GRAY_DARK)


def slide_14_funnel(prs):
    """슬라이드 14: 캠페인 퍼널"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 4 · Phase 5", "캠페인 퍼널 — 어디서 고객을 잃고 있는가?")

    # 퍼널 차트 이미지 삽입
    img_path = f"{CHART_DIR}/layer5_campaign_funnel.png"
    img_inserted = add_image(sl, img_path,
                              100000, 820000, SLIDE_W // 2 - 150000, SLIDE_H - 1_100_000)

    # 수치 및 인사이트 (우측)
    RIGHT_L = SLIDE_W // 2 + 100000
    RIGHT_W = SLIDE_W // 2 - 200000
    TOP = 820000

    metrics = [
        ("발송", "100%",    BLACK),
        ("Open Rate", "53.68%", ACCENT_GREEN),
        ("CTR", "~8%",     ACCENT_BLUE),
        ("CVR", "0.53%",   ACCENT_RED),
    ]
    for i, (label, val, color) in enumerate(metrics):
        top = TOP + i * 900000
        add_rect(sl, RIGHT_L, top, RIGHT_W, 750000,
                 fill_color=GRAY_LIGHT, line_color=color, line_width=Pt(2))
        add_text(sl, label, RIGHT_L + 80000, top + 80000, RIGHT_W//2, 280000,
                 font_size=14, bold=True, color=GRAY_DARK)
        add_text(sl, val, RIGHT_L + RIGHT_W//2, top + 60000, RIGHT_W//2 - 80000, 380000,
                 font_size=26, bold=True, color=color, align=PP_ALIGN.RIGHT)

    add_text(sl, "⚠ 병목: 클릭→구매 구간\n→ 상품 페이지 UX 또는 가격 저항",
             RIGHT_L, TOP + 3_700_000, RIGHT_W, 500000,
             font_size=14, bold=True, color=ACCENT_RED)


def slide_15_ab(prs):
    """슬라이드 15: A/B 테스트"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 4 · Phase 5", "A/B 테스트 — 어떤 메시지가 더 잘 팔리는가?")

    img_path = f"{CHART_DIR}/layer5_ab_test_result.png"
    add_image(sl, img_path, 100000, 820000, SLIDE_W // 2 - 150000, SLIDE_H - 1_100_000)

    RIGHT_L = SLIDE_W // 2 + 100000
    RIGHT_W = SLIDE_W // 2 - 200000

    add_text(sl, "검정 결과",
             RIGHT_L, 820000, RIGHT_W, 320000,
             font_size=18, bold=True, color=BLACK)

    add_text(sl, "χ² = 2002.6   |   p < 0.0001",
             RIGHT_L, 1_180_000, RIGHT_W, 320000,
             font_size=16, bold=True, color=ACCENT_RED)
    add_text(sl, "통계적으로 매우 유의미한 차이",
             RIGHT_L, 1_480_000, RIGHT_W, 280000,
             font_size=14, color=GRAY_DARK)

    results = [
        ("A  Ranking",  "CTR  15.16%", "인기순 추천", ACCENT_BLUE),
        ("B  Curation", "CVR  12.73%", "개인화 큐레이션", ACCENT_GREEN),
    ]
    for i, (group, metric, desc, color) in enumerate(results):
        top = 1_900_000 + i * 1_100_000
        add_rect(sl, RIGHT_L, top, RIGHT_W, 950000, fill_color=GRAY_LIGHT, line_color=color, line_width=Pt(2))
        add_text(sl, group, RIGHT_L + 80000, top + 80000, RIGHT_W - 160000, 280000,
                 font_size=15, bold=True, color=color)
        add_text(sl, metric, RIGHT_L + 80000, top + 380000, RIGHT_W - 160000, 320000,
                 font_size=20, bold=True, color=color)
        add_text(sl, desc, RIGHT_L + 80000, top + 680000, RIGHT_W - 160000, 220000,
                 font_size=13, color=GRAY_DARK)

    add_text(sl, "✅  CTR 목표 → Ranking   |   CVR·매출 목표 → Curation",
             RIGHT_L, SLIDE_H - 600000, RIGHT_W, 380000,
             font_size=13, bold=True, color=ACCENT_RED)


def slide_16_roas(prs):
    """슬라이드 16: ROAS 시뮬레이션"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 4 · Phase 5", "ROAS 시뮬레이션 — 캠페인 예산 승인 근거")

    img_path = f"{CHART_DIR}/layer5_roas_simulation.png"
    add_image(sl, img_path, 100000, 820000, SLIDE_W // 2 - 150000, SLIDE_H - 1_200_000)

    RIGHT_L = SLIDE_W // 2 + 100000
    RIGHT_W = SLIDE_W // 2 - 200000

    add_text(sl, "세그먼트별 ROAS",
             RIGHT_L, 820000, RIGHT_W, 320000,
             font_size=18, bold=True, color=BLACK)

    roas_data = [
        ("Champions",       "2,595배", KAKAO_YELLOW, BLACK),
        ("Loyal Customers", "890배",   ACCENT_GREEN, WHITE),
        ("At Risk",         "340배",   ACCENT_BLUE,  WHITE),
        ("전체 평균",        "137배",   GRAY_MID,     WHITE),
    ]
    for i, (seg, roas, bg, fg) in enumerate(roas_data):
        top = 1_200_000 + i * 750000
        add_rect(sl, RIGHT_L, top, RIGHT_W, 620000, fill_color=bg)
        add_text(sl, seg, RIGHT_L + 80000, top + 80000, RIGHT_W//2, 260000,
                 font_size=14, bold=True, color=fg)
        add_text(sl, roas, RIGHT_L + RIGHT_W//2, top + 60000, RIGHT_W//2 - 80000, 360000,
                 font_size=24, bold=True, color=fg, align=PP_ALIGN.RIGHT)

    add_text(sl, "감도 분석 (3 시나리오)",
             RIGHT_L, 4_300_000, RIGHT_W, 300000,
             font_size=15, bold=True, color=BLACK)
    add_text(sl, "낙관적 +20% CVR → ROAS 165배\n기본 시나리오        → ROAS 137배\n보수적 -20% CVR → ROAS 110배\n\n→ 어느 시나리오도 ROAS > 100배",
             RIGHT_L, 4_620_000, RIGHT_W, 850000,
             font_size=13, color=GRAY_DARK)


def slide_17_crm_plan(prs):
    """슬라이드 17: CRM 액션 플랜"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 4 · Phase 5", "CRM 액션 플랜 — 어떤 세그먼트에 무슨 메시지를 언제?")

    img_path = f"{CHART_DIR}/layer5_segment_message_matrix.png"
    inserted = add_image(sl, img_path, 100000, 820000, SLIDE_W // 2 - 150000, SLIDE_H - 1_100_000)

    RIGHT_L = SLIDE_W // 2 + 100000
    RIGHT_W = SLIDE_W // 2 - 200000

    add_text(sl, "3종 즉시 실행 캠페인",
             RIGHT_L, 820000, RIGHT_W, 320000,
             font_size=18, bold=True, color=BLACK)

    campaigns = [
        ("VIP 로열티",   "Champions",               "얼리 액세스 + 신상품 초대", "상시", KAKAO_YELLOW, BLACK),
        ("재활성화",     "At Risk / Can't Lose Them", "개인화 메시지 + 쿠폰",     "마지막 구매 D+60", ACCENT_RED, WHITE),
        ("Occasion 트리거", "전 세그먼트",            "생일/기념일 맞춤 큐레이션", "이벤트 D-30", ACCENT_BLUE, WHITE),
    ]
    for i, (name, target, msg, timing, bg, fg) in enumerate(campaigns):
        top = 1_220_000 + i * 1_450_000
        add_rect(sl, RIGHT_L, top, RIGHT_W, 1_300_000, fill_color=bg)
        add_text(sl, name, RIGHT_L + 80000, top + 80000, RIGHT_W - 160000, 280000,
                 font_size=16, bold=True, color=fg)
        add_text(sl, f"타겟: {target}", RIGHT_L + 80000, top + 380000, RIGHT_W - 160000, 240000,
                 font_size=13, color=fg)
        add_text(sl, f"메시지: {msg}", RIGHT_L + 80000, top + 600000, RIGHT_W - 160000, 240000,
                 font_size=13, color=fg)
        add_text(sl, f"타이밍: {timing}", RIGHT_L + 80000, top + 820000, RIGHT_W - 160000, 380000,
                 font_size=13, bold=True, color=fg)


def slide_18_seasonal_strategy(prs):
    """슬라이드 18: 시즌/비시즌 전략 분리"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "PART 4 · Phase 5", "연간 캠페인 전략 — 시즌 vs 비시즌 분리 운영")

    img_path = f"{CHART_DIR}/layer5_seasonal_vs_normal.png"
    add_image(sl, img_path, 100000, 820000, SLIDE_W // 2 - 150000, SLIDE_H - 1_100_000)

    RIGHT_L = SLIDE_W // 2 + 100000
    RIGHT_W = SLIDE_W // 2 - 200000

    strategies = [
        ("📅  시즌 캠페인",
         ["목표: 대량 도달, 높은 Open Rate",
          "메시지: Ranking (인기 상품)",
          "타이밍: 이벤트 D-7 ~ D-1",
          "효과: 광범위 신규 유입"],
         RGBColor(0xFF,0xA0,0x00)),
        ("🎯  비시즌 캠페인",
         ["목표: 전환율 극대화, CVR",
          "메시지: Curation (개인화)",
          "타이밍: Golden Time (수신 후 30일)",
          "효과: 기존 사용자 재활성화"],
         ACCENT_BLUE),
    ]
    for i, (title, items, color) in enumerate(strategies):
        top = 820000 + i * 2_500_000
        add_rect(sl, RIGHT_L, top, RIGHT_W, 2_300_000, fill_color=GRAY_LIGHT, line_color=color, line_width=Pt(2))
        add_rect(sl, RIGHT_L, top, RIGHT_W, 380000, fill_color=color)
        add_text(sl, title, RIGHT_L + 80000, top + 100000, RIGHT_W - 160000, 280000,
                 font_size=16, bold=True, color=WHITE)
        for j, item in enumerate(items):
            add_text(sl, f"• {item}",
                     RIGHT_L + 80000, top + 450000 + j * 420000, RIGHT_W - 160000, 350000,
                     font_size=13, color=GRAY_DARK)


def slide_19_summary(prs):
    """슬라이드 19: 핵심 발견 요약"""
    sl = blank_slide(prs)
    fill_bg(sl, BLACK)
    add_rect(sl, 0, 0, SLIDE_W, 700000, fill_color=KAKAO_YELLOW)
    add_text(sl, "핵심 발견 요약 — 5 Phase 비즈니스 임팩트",
             100000, 130000, SLIDE_W - 200000, 480000,
             font_size=24, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    findings = [
        ("Phase 1", "GMV +29.7% YoY,\nTrend 성분 독립 성장",    "비시즌 CRM 투자 정당화"),
        ("Phase 2", "Champions가 GMV 44.3%\n상위 2개 세그먼트 = 63%", "CRM 예산 집중 근거"),
        ("Phase 3", "M+11 LTV ₩203,169\n시즌=비시즌 품질 동일",   "CAC 상한 & 시즌 예산 확대"),
        ("Phase 4", "K=1.559, Golden Time 30일\n생일 Occasion 1위",  "바이럴 자동화 파이프라인"),
        ("Phase 5", "Champions ROAS 2,595배\n총 GMV ₩6.02억",      "3종 CRM 즉시 실행 가능"),
    ]

    BOX_W = (SLIDE_W - 200000) // 5 - 30000
    for i, (phase, finding, impact) in enumerate(findings):
        lft = 100000 + i * (BOX_W + 40000)
        # 상단 Phase 레이블
        add_rect(sl, lft, 800000, BOX_W, 380000, fill_color=KAKAO_YELLOW)
        add_text(sl, phase, lft + 40000, 840000, BOX_W - 80000, 300000,
                 font_size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        # 본문
        add_rect(sl, lft, 1_200_000, BOX_W, 2_800_000, fill_color=RGBColor(0x22,0x22,0x22))
        add_text(sl, finding, lft + 60000, 1_300_000, BOX_W - 120000, 1_400_000,
                 font_size=13, color=WHITE)
        add_rect(sl, lft, 4_050_000, BOX_W, 80000, fill_color=KAKAO_YELLOW)
        add_text(sl, impact, lft + 60000, 4_170_000, BOX_W - 120000, 800000,
                 font_size=12, bold=True, color=KAKAO_YELLOW)


def slide_20_next_steps(prs):
    """슬라이드 20: 다음 단계"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "CONCLUSION", "다음 단계 — 분석을 실행으로")

    steps = [
        ("🔧  Occasion 파이프라인",
         "생일/기념일 D-30 자동 트리거 시스템 구축\n수신 후 D+7/14/30 자동 넛지 시퀀스",
         "우선순위: 높음  |  난이도: 중간", ACCENT_BLUE),
        ("🔄  세그먼트 자동 업데이트",
         "월 1회 RFM 재계산 → 세그먼트 자동 갱신\nChampions 이탈 감지 알럿 시스템",
         "우선순위: 높음  |  난이도: 낮음", ACCENT_GREEN),
        ("🧪  멀티암 A/B 테스트",
         "현재 이분법(Ranking/Curation) →\nHybrid 멀티암 테스트로 고도화",
         "우선순위: 중간  |  난이도: 중간", RGBColor(0x77,0x17,0xAA)),
        ("📊  LTV 기반 CAC 차등화",
         "세그먼트별 예상 LTV 산출 →\n허용 CAC 상한선 차등 적용",
         "우선순위: 중간  |  난이도: 낮음", ACCENT_RED),
    ]

    BOX_W = (SLIDE_W - 300000) // 2
    for i, (title, body, meta, color) in enumerate(steps):
        col = i % 2
        row = i // 2
        lft = 100000 + col * (BOX_W + 100000)
        top = 900000 + row * 2_300_000
        add_rect(sl, lft, top, BOX_W, 2_100_000, fill_color=GRAY_LIGHT, line_color=color, line_width=Pt(2))
        add_rect(sl, lft, top, BOX_W, 350000, fill_color=color)
        add_text(sl, title, lft + 80000, top + 80000, BOX_W - 160000, 260000,
                 font_size=15, bold=True, color=WHITE)
        add_text(sl, body, lft + 80000, top + 420000, BOX_W - 160000, 1_100_000,
                 font_size=13, color=GRAY_DARK)
        add_text(sl, meta, lft + 80000, top + 1_700_000, BOX_W - 160000, 280000,
                 font_size=12, bold=True, color=color)


def slide_21_limitations(prs):
    """슬라이드 21: 분석 한계 & 개선 방향"""
    sl = blank_slide(prs)
    fill_bg(sl, WHITE)
    add_header_bar(sl, "CONCLUSION", "분석 한계 & 개선 방향 — 솔직한 자기 평가")

    add_text(sl, "실제 서비스라면 보완할 점을 명확히 인식하고 있습니다.",
             100000, 820000, SLIDE_W - 200000, 350000,
             font_size=16, color=GRAY_DARK)

    limitations = [
        ("데이터 가정",
         "합성 데이터 사용 → 실제 분포와 다를 수 있음\n특히 K-factor, Occasion 비율은 실데이터 검증 필요",
         "개선: 실제 로그 데이터 + A/B 실험군 분리"),
        ("인과관계",
         "상관관계 분석 위주 → 인과추론(Causal Inference) 미적용\n예: 'Occasion이 원인이다'는 아직 가설",
         "개선: DiD, IV, RDD 등 인과추론 기법 적용"),
        ("K-factor 측정",
         "추천 경로 추적 데이터 불완전\n순수 바이럴 vs 유기적 성장 분리 어려움",
         "개선: attribution 모델 도입, 링크 트래킹"),
        ("세그먼트 경계",
         "RFM 분위 경계값 임의 설정\n비즈니스 컨텍스트 반영한 튜닝 필요",
         "개선: 비즈니스팀과 협업해 경계값 재정의"),
    ]

    COL_W = (SLIDE_W - 300000) // 2
    for i, (title, prob, fix) in enumerate(limitations):
        col = i % 2
        row = i // 2
        lft = 100000 + col * (COL_W + 100000)
        top = 1_300_000 + row * 2_200_000
        add_rect(sl, lft, top, COL_W, 2_000_000,
                 fill_color=GRAY_LIGHT, line_color=GRAY_MID, line_width=Pt(1))
        add_text(sl, f"⚠ {title}", lft + 80000, top + 100000, COL_W - 160000, 280000,
                 font_size=14, bold=True, color=ACCENT_RED)
        add_text(sl, prob, lft + 80000, top + 430000, COL_W - 160000, 750000,
                 font_size=12, color=GRAY_DARK)
        add_rect(sl, lft + 80000, top + 1_250_000, COL_W - 160000, 50000,
                 fill_color=ACCENT_GREEN)
        add_text(sl, fix, lft + 80000, top + 1_330_000, COL_W - 160000, 550000,
                 font_size=12, bold=True, color=ACCENT_GREEN)


# ─── 메인 실행 ────────────────────────────────────────────

def main():
    prs = new_prs()

    print("슬라이드 생성 중...")
    slide_01_cover(prs)       ; print("  1/21 표지")
    slide_02_problem(prs)     ; print("  2/21 문제 정의")
    slide_03_design(prs)      ; print("  3/21 분석 설계")
    slide_04_gmv(prs)         ; print("  4/21 GMV 성장 구조")
    slide_05_season(prs)      ; print("  5/21 시즌 패턴")
    slide_06_category(prs)    ; print("  6/21 카테고리 믹스")
    slide_07_rfm(prs)         ; print("  7/21 RFM 세그멘테이션")
    slide_08_pareto(prs)      ; print("  8/21 GMV Pareto")
    slide_09_ltv(prs)         ; print("  9/21 LTV 커브")
    slide_10_cohort(prs)      ; print(" 10/21 시즌 코호트 품질")
    slide_11_kfactor(prs)     ; print(" 11/21 K-factor")
    slide_12_reciprocity(prs) ; print(" 12/21 상호성 & Golden Time")
    slide_13_occasion(prs)    ; print(" 13/21 Occasion 트리거")
    slide_14_funnel(prs)      ; print(" 14/21 캠페인 퍼널")
    slide_15_ab(prs)          ; print(" 15/21 A/B 테스트")
    slide_16_roas(prs)        ; print(" 16/21 ROAS 시뮬레이션")
    slide_17_crm_plan(prs)    ; print(" 17/21 CRM 액션 플랜")
    slide_18_seasonal_strategy(prs) ; print(" 18/21 시즌/비시즌 전략")
    slide_19_summary(prs)     ; print(" 19/21 핵심 발견 요약")
    slide_20_next_steps(prs)  ; print(" 20/21 다음 단계")
    slide_21_limitations(prs) ; print(" 21/21 분석 한계")

    prs.save(OUTPUT)
    print(f"\nDone: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()