# -*- coding: utf-8 -*-
import sys
from pptx import Presentation

pptx_path = "selfmade/PPT 관련/kakao_gift_crm_analysis_개쉐파크_v01.pptx"
prs = Presentation(pptx_path)
slide = prs.slides[14]  # index 14 = 슬라이드 15

fixes = {
    "Google Shape;910;p17": ("캠페인: 240,662건",   "캠페인: 210,716건"),
    "Google Shape;920;p17": ("240,662건   100%",    "210,716건   100%"),
    "Google Shape;922;p17": ("129,147건   53.68%",  "126,424건   60.00%"),
    "Google Shape;924;p17": ("13,906건   5.78%",    "13,285건    6.30%"),
    "Google Shape;926;p17": ("1,276건   0.53%",     "1,294건     0.61%"),
}

for shape in slide.shapes:
    if shape.name not in fixes or not shape.has_text_frame:
        continue
    old, new = fixes[shape.name]
    for para in shape.text_frame.paragraphs:
        if old in para.text:
            if para.runs:
                para.runs[0].text = new
                for run in para.runs[1:]:
                    run.text = ""
            sys.stdout.buffer.write(f"  수정: {repr(old)} -> {repr(new)}\n".encode("utf-8"))

prs.save(pptx_path)
sys.stdout.buffer.write("저장 완료\n".encode("utf-8"))

# 검증
prs2 = Presentation(pptx_path)
slide2 = prs2.slides[14]
sys.stdout.buffer.write("\n=== 수정 결과 확인 ===\n".encode("utf-8"))
for shape in slide2.shapes:
    if shape.name in fixes and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                sys.stdout.buffer.write(f"[{shape.name}] {para.text}\n".encode("utf-8"))
