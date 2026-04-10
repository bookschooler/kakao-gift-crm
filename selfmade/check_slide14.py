# -*- coding: utf-8 -*-
import sys
from pptx import Presentation

prs = Presentation("selfmade/PPT 관련/kakao_gift_crm_analysis_개쉐파크_v01.pptx")

# PDF Page 14 = PPT index 14 (slide 15) 확인
slide = prs.slides[14]
sys.stdout.buffer.write("=== PPT Slide 15 (index 14) 전체 shape ===\n".encode("utf-8"))
for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                line = f"[Shape {i} | {shape.name}] {repr(t)}\n"
                sys.stdout.buffer.write(line.encode("utf-8"))
