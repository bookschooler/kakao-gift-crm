# -*- coding: utf-8 -*-
"""
PPT 수치 수정 스크립트
======================
슬라이드 15, 17, 19의 구버전 수치를 새 데이터 기반 수치로 교체.

수정 내용:
  Slide 15: 240,662 → 210,716 (캠페인 퍼널 발송 건수)
  Slide 17: ROAS 구 수치 → 신 수치 (Can't Lose Them 96x 등)
  Slide 19: 240,662 → 210,716 (요약 슬라이드)
"""
import sys
from pptx import Presentation

pptx_path = "selfmade/PPT 관련/kakao_gift_crm_analysis_개쉐파크_v01.pptx"
prs = Presentation(pptx_path)

# ─────────────────────────────────────────────
# 슬라이드별 수정 맵: {slide_index: [(old_text, new_text), ...]}
# ─────────────────────────────────────────────
SLIDE_FIXES = {
    # 슬라이드 15 (index 14): 퍼널 발송 건수
    14: [
        ("240,662건   100%",          "210,716건   100%"),
        ("캠페인: 240,662건",          "캠페인: 210,716건"),
        # 퍼널 수치도 실측으로 교체
        ("129,147건   53.68%",        "126,424건   60.00%"),
        ("13,906건   5.78%",          "13,285건    6.30%"),
        ("1,276건   0.53%",           "1,294건     0.61%"),
    ],
    # 슬라이드 17 (index 16): ROAS 시뮬레이션
    16: [
        # 핵심 지표 카드
        ("세그먼트별 ROAS (K-factor 1.559 반영, 비용 15원/건)",
         "세그먼트별 ROAS (K-factor 1.559 반영, 비용 15원/건)"),  # 유지
        ("모든 시나리오에서 ROAS > 100배를 달성합니다.",
         "모든 시나리오에서 ROAS > 30배를 달성합니다."),
        # 구 수치 → 신 수치
        ("Champions ROAS",            "Champions ROAS"),       # 유지
        ("2,595x ↑12%",              "82x  ↑2.0x vs 평균"),
        ("At Risk ROAS",              "Can't Lose Them ROAS"),
        ("1,261x ↑8%",               "96x  ↑2.4x vs 평균"),
        ("Need Attention",            "Loyal Customers"),
        ("1,252x ↑5%",               "65x  ↑1.6x vs 평균"),
        ("Loyal ROAS",                "전체 평균 ROAS"),
        ("1,186x ↑3%",               "40x  (15원/건 기준)"),
        # 실행 제언 표
        ("2,595x",                    "82x"),
        ("1,261x",                    "96x"),
        ("1,500x",                    "30x (보수적)"),
        # 텍스트 설명
        ("비용 변화에도 순위 불변 → 안정적 수익성 확보 가능",
         "비용 변화에도 순위 불변 - 안정적 수익성 확보 가능"),
    ],
    # 슬라이드 19 (index 18): 요약
    18: [
        ("총 캠페인   240,662건",      "총 캠페인   210,716건"),
        ("240,662건",                  "210,716건"),
    ],
}


def replace_text_in_shape(shape, old, new):
    """shape 내 텍스트를 run 단위로 교체. 반환: 교체 성공 여부."""
    if not shape.has_text_frame:
        return False
    changed = False
    tf = shape.text_frame
    for para in tf.paragraphs:
        # 단락 전체 텍스트에 old가 포함되면
        full = para.text
        if old in full:
            # run이 1개면 바로 교체
            if len(para.runs) == 1:
                para.runs[0].text = para.runs[0].text.replace(old, new)
                changed = True
            elif len(para.runs) == 0:
                pass  # 빈 단락
            else:
                # run이 여러 개 → 첫 번째 run에 전체 텍스트 몰아넣기
                new_full = full.replace(old, new)
                para.runs[0].text = new_full
                for run in para.runs[1:]:
                    run.text = ""
                changed = True
    return changed


total_changes = 0
for slide_idx, fixes in SLIDE_FIXES.items():
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1
    for old_text, new_text in fixes:
        if old_text == new_text:
            continue  # 변경 없는 항목 스킵
        for shape in slide.shapes:
            if replace_text_in_shape(shape, old_text, new_text):
                msg = f"  [Slide {slide_num}] '{old_text}' -> '{new_text}'"
                sys.stdout.buffer.write(msg.encode("utf-8") + b"\n")
                total_changes += 1

prs.save(pptx_path)
sys.stdout.buffer.write(
    f"\n저장 완료: {pptx_path}\n총 {total_changes}건 교체\n".encode("utf-8")
)

# ─── 검증 ─────────────────────────────────────
sys.stdout.buffer.write(b"\n=== 수정 결과 검증 ===\n")
prs2 = Presentation(pptx_path)
for slide_idx in SLIDE_FIXES:
    slide = prs2.slides[slide_idx]
    sys.stdout.buffer.write(f"\n[Slide {slide_idx+1}]\n".encode("utf-8"))
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t and any(kw in t for kw in
                             ["건", "ROAS", "x ", "x\u3000", "96", "82", "40", "30", "평균", "보수"]):
                    sys.stdout.buffer.write(f"  {t}\n".encode("utf-8"))