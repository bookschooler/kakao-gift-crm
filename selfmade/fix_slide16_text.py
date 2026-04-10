# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt

pptx_path = "selfmade/PPT 관련/kakao_gift_crm_analysis_개쉐파크_v01.pptx"
prs = Presentation(pptx_path)

slide = prs.slides[15]

shape19 = None
shape20 = None
for shape in slide.shapes:
    if shape.name == "Google Shape;980;p18":
        shape19 = shape
    if shape.name == "Google Shape;981;p18":
        shape20 = shape

# Shape 20: 타이틀 수정
tf20 = shape20.text_frame
tf20.paragraphs[0].runs[0].text = "Curation 기본 전략, Ranking 제한 사용:"

# Shape 19: 실측 수치 기반 실행 제언 3줄
# 실측: CTR(발송기준) Ranking 4.43% vs Curation 6.11%
#        CVR(발송기준) Ranking 0.22% vs Curation 0.71%
#        Block Rate    Ranking 1.19% vs Curation 0.69%
new_lines = [
    " • 전환 목표 → Curation 기본 집행: CTR 6.11% vs Ranking 4.43%, CVR 0.71% vs 0.22% — 모든 KPI에서 우세. Champions·Can't Lose Them 고AOV 세그먼트에 집중",
    " • 이탈 방지 → Curation 우선: Block Rate 0.69%(Ranking 1.19% 대비 절반). At Risk 재활성화 캠페인에서 관계 훼손 없이 재접근 가능",
    " • Ranking 제한 사용: 신규·잠재 고객 대상 노출 확장 목적에만 한정. 전환 목적 캠페인에서는 Curation으로 전환 권장",
]

tf19 = shape19.text_frame
for i, para in enumerate(tf19.paragraphs):
    if i < len(new_lines):
        for run in para.runs:
            run.text = ""
        if para.runs:
            para.runs[0].text = new_lines[i]

prs.save(pptx_path)
print("저장 완료")
print("=== 수정된 실행 제언 ===")
for shape in slide.shapes:
    if shape.has_text_frame and shape.name in [
        "Google Shape;979;p18",
        "Google Shape;980;p18",
        "Google Shape;981;p18",
    ]:
        print(f"\n[{shape.name}]")
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                print(f"  {para.text}")
