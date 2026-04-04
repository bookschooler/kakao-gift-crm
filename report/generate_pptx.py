"""
카카오톡 선물하기 CRM 분석 — PPTX 생성 스크립트 (13슬라이드)
ANALYSIS_RESULTS.md 수치 기반 / python-pptx / McKinsey-BCG 스타일
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE = r"C:\Users\user\Desktop\pjt\portfolio\kakao_gift"
CHART_DIR = os.path.join(BASE, "report", "charts")
OUT_PATH = os.path.join(BASE, "report", "kakao_gift_analysis.pptx")

# ── 색상 ───────────────────────────────────────────
PRIMARY   = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT    = RGBColor(0xFE, 0xE5, 0x00)
HIGHLIGHT = RGBColor(0x22, 0x51, 0xFF)
SUCCESS   = RGBColor(0x29, 0xBA, 0x74)
DANGER    = RGBColor(0xFF, 0x5A, 0x5F)
BG        = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x6B, 0x72, 0x80)
LIGHT     = RGBColor(0xF3, 0xF4, 0xF6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ── 슬라이드 크기: 와이드스크린 13.33 × 7.5 ────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃

# ─────────────────────────────────────────────────────
# 헬퍼 함수
# ─────────────────────────────────────────────────────
def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
             word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run.font.name = '맑은 고딕'
    return txBox

def add_picture_safe(slide, img_path, left, top, width, height):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)
        return True
    return False

def top_bar(slide, slide_num, section=""):
    """상단 컬러바 + 슬라이드 번호"""
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.15), fill_color=PRIMARY)
    add_text(slide, f"{slide_num:02d}", W - Inches(0.7), Inches(0.03),
             Inches(0.6), Inches(0.2), font_size=9, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

def footer(slide, note="출처: 카카오 선물하기 포트폴리오 프로젝트 | 2023-2024 | 모의 데이터"):
    add_text(slide, note, Inches(0.3), H - Inches(0.25),
             Inches(12), Inches(0.2), font_size=7, color=MUTED)

def section_title(slide, title, font_size=28):
    add_text(slide, title, Inches(0.4), Inches(0.22), Inches(12.0), Inches(0.55),
             font_size=font_size, bold=True, color=PRIMARY)

def takeaway_box(slide, text, top=Inches(0.85)):
    box = add_shape(slide, Inches(0.4), top, Inches(12.5), Inches(0.42), fill_color=ACCENT)
    add_text(slide, "  " + text, Inches(0.4), top, Inches(12.5), Inches(0.42),
             font_size=11, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────
# SLIDE 01 — 표지
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
# 전체 배경 (딥 네이비)
add_shape(slide, Inches(0), Inches(0), W, H, fill_color=PRIMARY)
# 우측 포인트 바
add_shape(slide, W - Inches(0.5), Inches(0), Inches(0.5), H, fill_color=ACCENT)
# 메인 제목
add_text(slide,
         "카카오톡 선물하기",
         Inches(0.8), Inches(1.8), Inches(10), Inches(1.0),
         font_size=44, bold=True, color=WHITE)
add_text(slide,
         "CRM 분석 보고서",
         Inches(0.8), Inches(2.7), Inches(10), Inches(0.9),
         font_size=40, bold=True, color=ACCENT)
# 부제
add_text(slide,
         "RFM 세그멘테이션 · LTV 코호트 · Viral Loop 분석 → CRM 액션 플랜",
         Inches(0.8), Inches(3.7), Inches(11), Inches(0.5),
         font_size=16, color=LIGHT)
# 구분선
add_shape(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.03), fill_color=ACCENT)
# 메타정보
add_text(slide,
         "분석 기간: 2023-01 ~ 2024-12   |   데이터: 57만 행, 5개 테이블",
         Inches(0.8), Inches(4.5), Inches(11), Inches(0.35),
         font_size=13, color=MUTED)
add_text(slide,
         "2026-04-04   |   Portfolio Project by Soyoung Yun",
         Inches(0.8), Inches(4.95), Inches(11), Inches(0.35),
         font_size=13, color=MUTED)
# KPI 카드 3개
kpi_data = [("₩137.2억", "총 GMV"), ("+31%", "YoY 성장"), ("K-factor\n3.95", "바이럴 성장")]
for i, (val, label) in enumerate(kpi_data):
    x = Inches(0.8 + i * 3.5)
    add_shape(slide, x, Inches(5.7), Inches(3.0), Inches(1.3),
              fill_color=RGBColor(0x2A, 0x2A, 0x4E))
    add_text(slide, val, x + Inches(0.15), Inches(5.78), Inches(2.7), Inches(0.65),
             font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.15), Inches(6.35), Inches(2.7), Inches(0.35),
             font_size=10, color=LIGHT, align=PP_ALIGN.CENTER)
print("슬라이드 01 완료")

# ─────────────────────────────────────────────────────
# SLIDE 02 — 목차
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 2)
section_title(slide, "분석은 5개 Layer로 구성 — GMV부터 캠페인 ROAS까지")
takeaway_box(slide, "핵심 KPI: 총 GMV ₩137.2억 | K-factor 3.95 | Champions ROAS 1,665배")

layers = [
    ("Layer 1", "EDA & 시즌 분석",      "GMV 트렌드 + YoY +31% + 카테고리 믹스"),
    ("Layer 2", "RFM 세그멘테이션",     "11개 세그먼트 | Champions 3.2% = GMV 9.3%"),
    ("Layer 3", "LTV 코호트 분석",      "1개월 Retention 15.1% | 12개월 LTV 2.7배"),
    ("Layer 4", "Viral Loop 분석",      "K-factor 3.95 | Reciprocity 67.6%"),
    ("Layer 5", "CRM 캠페인 전략",      "세그먼트별 ROAS | 3대 액션 플랜"),
]
colors = [HIGHLIGHT, SUCCESS, ACCENT, DANGER, PRIMARY]
for i, (layer, title, desc) in enumerate(layers):
    y = Inches(1.5 + i * 1.06)
    add_shape(slide, Inches(0.4), y, Inches(0.08), Inches(0.75), fill_color=colors[i])
    add_text(slide, layer, Inches(0.6), y + Inches(0.05), Inches(1.3), Inches(0.35),
             font_size=9, bold=True, color=MUTED)
    add_text(slide, title, Inches(0.6), y + Inches(0.32), Inches(3.5), Inches(0.4),
             font_size=14, bold=True, color=PRIMARY)
    add_text(slide, desc, Inches(4.3), y + Inches(0.25), Inches(8.5), Inches(0.45),
             font_size=11, color=MUTED)

# KPI 미리보기 박스
kpi_items = [("₩137.2억", "총 GMV (2년)", SUCCESS),
             ("K=3.95", "바이럴 K-factor", HIGHLIGHT),
             ("1,665x", "Champions ROAS", ACCENT)]
for i, (val, lbl, col) in enumerate(kpi_items):
    x = Inches(9.8)
    y = Inches(1.55 + i * 1.6)
    add_shape(slide, x, y, Inches(3.1), Inches(1.2), fill_color=LIGHT)
    add_shape(slide, x, y, Inches(0.07), Inches(1.2), fill_color=col)
    add_text(slide, val, x + Inches(0.2), y + Inches(0.1), Inches(2.7), Inches(0.6),
             font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, lbl, x + Inches(0.1), y + Inches(0.75), Inches(2.9), Inches(0.35),
             font_size=10, color=MUTED, align=PP_ALIGN.CENTER)
footer(slide)
print("슬라이드 02 완료")

# ─────────────────────────────────────────────────────
# SLIDE 03 — 데이터 개요
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 3)
section_title(slide, "57만 행, 5개 테이블, 2년 분석 — 선물하기 고유 구조 반영")
takeaway_box(slide, "발신자(구매자) ≠ 수신자 구조 + Viral Loop 추적이 이 분석의 핵심 차별점")

tables_data = [
    ("users", "50,000명", "유저 프로필"),
    ("orders", "200,011건", "발신 기준 주문"),
    ("gift_receipts", "200,011건", "수신 기록 (1:1)"),
    ("campaigns", "88건", "CRM 캠페인 기획"),
    ("campaign_logs", "339,233건", "오픈/클릭/구매 로그"),
]
t_colors = [HIGHLIGHT, SUCCESS, ACCENT, PRIMARY, DANGER]
for i, (name, count, desc) in enumerate(tables_data):
    x = Inches(0.4 + i * 2.55)
    add_shape(slide, x, Inches(1.5), Inches(2.3), Inches(3.5), fill_color=LIGHT)
    add_shape(slide, x, Inches(1.5), Inches(2.3), Inches(0.08), fill_color=t_colors[i])
    add_text(slide, name, x + Inches(0.2), Inches(1.6), Inches(2.0), Inches(0.4),
             font_size=13, bold=True, color=t_colors[i])
    add_text(slide, count, x + Inches(0.2), Inches(2.05), Inches(2.0), Inches(0.6),
             font_size=20, bold=True, color=PRIMARY)
    add_text(slide, desc, x + Inches(0.2), Inches(2.65), Inches(2.0), Inches(0.35),
             font_size=10, color=MUTED)

# 분석 스펙
specs = [
    ("분석 기간", "2023-01-01 ~ 2024-12-31"),
    ("총 행 수", "약 57만 행"),
    ("RFM 윈도우", "최근 12개월"),
    ("코호트 수", "24개 (월별)"),
    ("캠페인 비용 가정", "15원/건"),
]
for i, (k, v) in enumerate(specs):
    y = Inches(5.15)
    x = Inches(0.4 + i * 2.55)
    add_shape(slide, x, y, Inches(2.3), Inches(0.75), fill_color=RGBColor(0xE8, 0xEA, 0xF0))
    add_text(slide, k, x + Inches(0.1), y + Inches(0.05), Inches(2.1), Inches(0.28),
             font_size=9, color=MUTED)
    add_text(slide, v, x + Inches(0.1), y + Inches(0.33), Inches(2.1), Inches(0.35),
             font_size=11, bold=True, color=PRIMARY)

add_text(slide, "분석 환경: Python (pandas, scipy, matplotlib) + BigQuery",
         Inches(0.4), Inches(6.1), Inches(12), Inches(0.3),
         font_size=10, color=MUTED)
footer(slide)
print("슬라이드 03 완료")

# ─────────────────────────────────────────────────────
# SLIDE 04 — Layer 1: GMV 트렌드
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 4, "Layer 1")
section_title(slide, "2024년 GMV ₩137.2억, YoY +31% — 1월 시즌이 연중 최고점")
takeaway_box(slide, "2024-01 GMV ₩10.2억 (피크) vs 2023-06 ₩2.9억 (저점) — 계절성이 사업 리듬을 결정")
add_picture_safe(slide, os.path.join(CHART_DIR, "chart_01_gmv_trend.png"),
                 Inches(0.4), Inches(1.45), Inches(12.5), Inches(5.3))
footer(slide, "주: 월별 GMV = 해당 월 완료 주문 기준 | YoY = 2024년 12개월 평균 vs 2023년")
print("슬라이드 04 완료")

# ─────────────────────────────────────────────────────
# SLIDE 05 — Layer 1: 시즌 분석 + 카테고리 믹스
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 5, "Layer 1")
section_title(slide, "빼빼로데이 ×12 폭증, beauty GMV 1위 — 시즌·카테고리 집중도 극심")
takeaway_box(slide, "시즌 이벤트 Top3(빼빼로·발렌타인·어버이날) 집중 운영 + voucher 단가 업셀링 병행 필요")

# 왼쪽: 시즌 이벤트 테이블
events = [
    ("빼빼로데이 11/11", "×12.0", ACCENT),
    ("어버이날 5/8", "×3.5", SUCCESS),
    ("설날·발렌타인 1~2월", "×3.0", HIGHLIGHT),
    ("스승의날·크리스마스", "×2.5", MUTED),
    ("블랙프라이데이", "×1.8", MUTED),
]
add_text(slide, "시즌 GMV 부스트 배율", Inches(0.4), Inches(1.45), Inches(5), Inches(0.35),
         font_size=12, bold=True, color=PRIMARY)
for i, (evt, boost, col) in enumerate(events):
    y = Inches(1.85 + i * 0.82)
    add_shape(slide, Inches(0.4), y, Inches(5.5), Inches(0.65), fill_color=LIGHT)
    add_shape(slide, Inches(0.4), y, Inches(0.07), Inches(0.65), fill_color=col)
    add_text(slide, evt, Inches(0.6), y + Inches(0.15), Inches(3.8), Inches(0.35),
             font_size=11, color=PRIMARY)
    add_text(slide, boost, Inches(5.1), y + Inches(0.12), Inches(0.9), Inches(0.42),
             font_size=18, bold=True, color=col, align=PP_ALIGN.RIGHT)

# 오른쪽: 카테고리 차트
add_picture_safe(slide, os.path.join(CHART_DIR, "chart_02_category_mix.png"),
                 Inches(6.2), Inches(1.4), Inches(6.9), Inches(4.8))

# 유저 분포
add_text(slide, "유저 분포: 여성 61.8% | 30대 37.9% | gift_received 채널 1위(40%)",
         Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.3),
         font_size=10, color=MUTED)
footer(slide)
print("슬라이드 05 완료")

# ─────────────────────────────────────────────────────
# SLIDE 06 — Layer 2: RFM 세그먼트 맵
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 6, "Layer 2")
section_title(slide, "11개 RFM 세그먼트 — 유저 분포와 GMV 기여도가 역전되는 세그먼트 존재")
takeaway_box(slide, "Hibernating이 유저 23.2%지만 GMV 11.7%, Champions는 유저 3.2%로 GMV 9.3% 창출")
add_picture_safe(slide, os.path.join(CHART_DIR, "chart_03_rfm_segments.png"),
                 Inches(0.3), Inches(1.45), Inches(12.7), Inches(5.3))
footer(slide, "주: RFM 분석 기간 2024-01-01~2024-12-31 | 분석 대상 44,713명")
print("슬라이드 06 완료")

# ─────────────────────────────────────────────────────
# SLIDE 07 — Layer 2: Champions 심층
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 7, "Layer 2")
section_title(slide, "Champions 1,438명 — 평균 구매액 ₩499K, GMV 기여 9.3%의 최고가치 그룹")
takeaway_box(slide, "Champions 이탈 1명 = ₩499,416 GMV 손실 → VIP 선제 관리가 최우선")

kpi_champ = [
    ("1,438명", "Champions 유저 수", "전체 3.2%"),
    ("₩499,416", "평균 구매액(M)", "Loyal 대비 +119%"),
    ("₩7.2억", "총 GMV 기여", "전체 GMV의 9.3%"),
    ("ROAS\n1,665배", "캠페인 효율", "비용 15원/건"),
]
for i, (val, label, note) in enumerate(kpi_champ):
    x = Inches(0.4 + i * 3.2)
    add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(2.2), fill_color=PRIMARY)
    add_text(slide, val, x + Inches(0.15), Inches(1.6), Inches(2.6), Inches(0.9),
             font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.1), Inches(2.55), Inches(2.7), Inches(0.38),
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, note, x + Inches(0.1), Inches(2.95), Inches(2.7), Inches(0.3),
             font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

# 세그먼트 비교 테이블
seg_compare = [
    ("세그먼트", "유저 수", "GMV 비중", "평균 M", ""),
    ("Champions", "1,438명", "9.3%", "₩499,416", ACCENT),
    ("Loyal", "8,708명", "25.8%", "₩228,300", SUCCESS),
    ("At Risk", "3,755명", "11.8%", "₩242,600", DANGER),
    ("Need Attention", "5,818명", "18.2%", "₩240,975", RGBColor(0xF5, 0x9E, 0x0B)),
]
for i, row in enumerate(seg_compare):
    y = Inches(3.9 + i * 0.52)
    bg = row[4] if len(row) > 4 and i > 0 else LIGHT
    if i == 0:
        bg = PRIMARY
    add_shape(slide, Inches(0.4), y, Inches(12.5), Inches(0.48), fill_color=bg)
    for j, (cell, width) in enumerate(zip(row[:4], [3.5, 2.2, 2.2, 2.5])):
        x = Inches(0.5 + sum([3.5, 2.2, 2.2, 2.5][:j]))
        text_color = WHITE if i <= 1 else PRIMARY
        add_text(slide, cell, Inches(0.5 + [0, 3.5, 5.7, 7.9][j]), y + Inches(0.1),
                 Inches(width - 0.1), Inches(0.32),
                 font_size=10, bold=(i == 0), color=WHITE if i == 0 else PRIMARY)

footer(slide)
print("슬라이드 07 완료")

# ─────────────────────────────────────────────────────
# SLIDE 08 — Layer 3: LTV 코호트 히트맵
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 8, "Layer 3")
section_title(slide, "12개월 누적 LTV ₩203,169 — 첫 달 대비 2.7배, 시즌 스파이크 뚜렷")
takeaway_box(slide, "1개월 Retention 15.1%지만 이탈 유저도 시즌 때 귀환 → 장기 LTV 관리가 단기 전환보다 중요")
add_picture_safe(slide, os.path.join(CHART_DIR, "chart_04_cohort_heatmap.png"),
                 Inches(0.3), Inches(1.45), Inches(12.7), Inches(5.3))
footer(slide, "주: 코호트 = 첫 구매월 기준 | LTV = 1인당 평균 누적 구매액")
print("슬라이드 08 완료")

# ─────────────────────────────────────────────────────
# SLIDE 09 — Layer 4: Viral Loop (K-factor)
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 9, "Layer 4")
section_title(slide, "K-factor 3.95 — 선물 1건이 평균 4명의 신규 구매자를 만든다")
takeaway_box(slide, "선물 수신이 신규 유저 유입 채널 1위(40%) — 바이럴이 유일한 무비용 성장 엔진")

# K-factor 시각화
add_shape(slide, Inches(0.4), Inches(1.45), Inches(4.5), Inches(5.35),
          fill_color=PRIMARY)
add_text(slide, "K-factor", Inches(0.6), Inches(1.65), Inches(4.0), Inches(0.5),
         font_size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, "3.95", Inches(0.6), Inches(2.15), Inches(4.0), Inches(1.2),
         font_size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "K > 1 = 바이럴 성장 확정", Inches(0.6), Inches(3.45), Inches(4.0), Inches(0.4),
         font_size=12, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)
add_text(slide, "Bootstrap CI: [3.87, 4.02]\n1,000회 반복 검증",
         Inches(0.6), Inches(3.95), Inches(4.0), Inches(0.6),
         font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

# K-factor 계산식
add_text(slide, "K = 발신자당 초대 수 × 전환율",
         Inches(0.6), Inches(4.7), Inches(4.0), Inches(0.35),
         font_size=10, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, "= 4.02명 × 98.1% = 3.948",
         Inches(0.6), Inches(5.05), Inches(4.0), Inches(0.35),
         font_size=10, color=ACCENT, align=PP_ALIGN.CENTER)

# Reciprocity
add_picture_safe(slide, os.path.join(CHART_DIR, "chart_05_viral_conversion.png"),
                 Inches(5.1), Inches(1.45), Inches(8.0), Inches(5.35))

# Reciprocity 박스
rec_data = [("30일 내 재발신", "58.4%"), ("60일 내", "63.3%"), ("90일 내", "67.6%"), ("최종 전환", "98.1%")]
for i, (label, val) in enumerate(rec_data):
    x = Inches(5.2 + i * 1.95)
    add_shape(slide, x, Inches(6.2), Inches(1.75), Inches(0.55), fill_color=LIGHT)
    add_text(slide, val, x, Inches(6.22), Inches(1.75), Inches(0.35),
             font_size=14, bold=True, color=HIGHLIGHT, align=PP_ALIGN.CENTER)

footer(slide)
print("슬라이드 09 완료")

# ─────────────────────────────────────────────────────
# SLIDE 10 — Layer 4: 바이럴 세대
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 10, "Layer 4")
section_title(slide, "Gen 1(바이럴 1세대) GMV ₩61.7억 > Gen 0(원조) ₩47.8억 — 바이럴 유저 품질이 더 높다")
takeaway_box(slide, "선물을 받은 후 보내는 사람이 원래 고객보다 더 많이, 더 자주 구매")

gen_data = [
    ("Gen 0", "최초 유입", "17,448명", "₩47.8억", "35.7%"),
    ("Gen 1", "1차 바이럴", "22,444명", "₩61.7억", "46.1%"),
    ("Gen 2", "2차 바이럴", "8,581명", "₩23.5억", "17.6%"),
    ("Gen 3+", "3차+", "1,527명", "₩4.3억", "3.2%"),
]
gen_colors = [HIGHLIGHT, ACCENT, SUCCESS, MUTED]
bar_max = 61.7
for i, (gen, desc, users, gmv, gmv_pct) in enumerate(gen_data):
    y = Inches(1.55 + i * 1.3)
    bar_w = Inches(12.3 * float(gmv.replace('₩','').replace('억','')) / bar_max)
    add_shape(slide, Inches(0.4), y, Inches(12.3), Inches(1.05), fill_color=LIGHT)
    add_shape(slide, Inches(0.4), y, bar_w, Inches(1.05), fill_color=gen_colors[i])
    add_text(slide, gen, Inches(0.55), y + Inches(0.08), Inches(1.2), Inches(0.45),
             font_size=14, bold=True, color=WHITE if i < 3 else PRIMARY)
    add_text(slide, desc, Inches(0.55), y + Inches(0.55), Inches(1.8), Inches(0.35),
             font_size=9, color=WHITE if i < 2 else PRIMARY)
    add_text(slide, f"유저 {users}", Inches(2.4), y + Inches(0.08), Inches(2.5), Inches(0.45),
             font_size=12, color=PRIMARY if i > 1 else WHITE)
    add_text(slide, f"GMV {gmv} ({gmv_pct})", Inches(5.0), y + Inches(0.08), Inches(3.5), Inches(0.45),
             font_size=13, bold=True, color=PRIMARY if i > 1 else WHITE)

footer(slide, "주: Referral Generation = gift_received 채널 기준 바이럴 세대 추적")
print("슬라이드 10 완료")

# ─────────────────────────────────────────────────────
# SLIDE 11 — Layer 5: 캠페인 성과
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 11, "Layer 5")
section_title(slide, "캠페인 Open Rate 53.7% 우수 — Purchase CVR 0.53%가 최적화 포인트")
takeaway_box(slide, "at_risk CVR 0.62%, promising CVR 0.81% — 넛지 대상 세그먼트에서 전환 효율 높음")

# 전체 평균 KPI
kpi_camp = [
    ("53.7%", "Open Rate", SUCCESS),
    ("5.8%", "Click Rate", HIGHLIGHT),
    ("1.0%", "Block Rate", DANGER),
    ("0.53%", "Purchase CVR", ACCENT),
]
for i, (val, label, col) in enumerate(kpi_camp):
    x = Inches(0.4 + i * 3.2)
    add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(1.3), fill_color=LIGHT)
    add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(0.07), fill_color=col)
    add_text(slide, val, x + Inches(0.15), Inches(1.58), Inches(2.6), Inches(0.65),
             font_size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.1), Inches(2.25), Inches(2.7), Inches(0.3),
             font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

# 세그먼트별 성과 테이블
seg_perf = [
    ("세그먼트", "Open Rate", "Click Rate", "Purchase CVR"),
    ("at_risk", "55.9%", "5.6%", "0.62%"),
    ("promising", "55.2%", "5.8%", "0.81%"),
    ("champions", "52.6%", "4.7%", "0.47%"),
    ("need_attention", "52.3%", "8.0%", "0.48%"),
    ("new_customers", "53.9%", "5.4%", "0.55%"),
    ("loyal_customers", "52.1%", "6.1%", "0.44%"),
]
for i, row in enumerate(seg_perf):
    y = Inches(3.0 + i * 0.55)
    bg = PRIMARY if i == 0 else (LIGHT if i % 2 == 1 else WHITE)
    add_shape(slide, Inches(0.4), y, Inches(12.5), Inches(0.52), fill_color=bg)
    for j, (cell, w) in enumerate(zip(row, [3.5, 3.0, 3.0, 3.0])):
        tc = WHITE if i == 0 else PRIMARY
        # CVR 강조
        if j == 3 and i > 0:
            val = float(cell.replace('%', ''))
            tc = SUCCESS if val >= 0.7 else (HIGHLIGHT if val >= 0.55 else PRIMARY)
        add_text(slide, cell, Inches(0.5 + [0, 3.5, 6.5, 9.5][j]), y + Inches(0.12),
                 Inches(w - 0.1), Inches(0.32),
                 font_size=10, bold=(i == 0 or (j == 3 and i > 0)),
                 color=tc)

footer(slide)
print("슬라이드 11 완료")

# ─────────────────────────────────────────────────────
# SLIDE 12 — Layer 5: CRM 액션 플랜
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 12, "Layer 5")
section_title(slide, "세그먼트별 ROAS 최대 1,665배 — 3개 우선 세그먼트로 GMV 최대화")
takeaway_box(slide, "Champions 보호 + At Risk 이탈 방지 + Need Attention 윈백 = 3대 CRM 전략")

add_picture_safe(slide, os.path.join(CHART_DIR, "chart_06_roas_bar.png"),
                 Inches(0.3), Inches(1.45), Inches(7.5), Inches(4.0))

# 3대 액션 플랜 카드
plans = [
    ("1. Champions", "ROAS 1,665x", "VIP 얼리버드 쿠폰\n시즌 선착순 혜택 선제 발송\n카카오 친구톡 (고도 개인화)", ACCENT, "3,591만원"),
    ("2. At Risk", "ROAS 809x", "감성 넛지 + 관계 기반 추천\n'소중한 분들 챙기세요'\n이탈 방지 최우선", DANGER, "4,555만원"),
    ("3. Need Attention", "ROAS 803x", "시즌 맞춤 윈백 캠페인\n할인 쿠폰 + 알림톡(비용 절감)\nGMV 2위 세그먼트 유지", SUCCESS, "7,010만원"),
]
for i, (title, roas, action, col, rev) in enumerate(plans):
    x = Inches(7.9)
    y = Inches(1.45 + i * 2.0)
    add_shape(slide, x, y, Inches(5.1), Inches(1.8), fill_color=LIGHT)
    add_shape(slide, x, y, Inches(0.07), Inches(1.8), fill_color=col)
    add_text(slide, title, x + Inches(0.2), y + Inches(0.08), Inches(3.0), Inches(0.38),
             font_size=13, bold=True, color=PRIMARY)
    add_text(slide, roas, x + Inches(3.3), y + Inches(0.08), Inches(1.6), Inches(0.38),
             font_size=13, bold=True, color=col, align=PP_ALIGN.RIGHT)
    add_text(slide, action, x + Inches(0.2), y + Inches(0.5), Inches(3.5), Inches(0.8),
             font_size=9, color=MUTED)
    add_text(slide, f"예상 수익 {rev}", x + Inches(0.2), y + Inches(1.45), Inches(3.0), Inches(0.28),
             font_size=9, bold=True, color=col)

footer(slide, "주: ROAS = (평균구매액 × CVR) / 15원 | 예상 수익 = 세그먼트 전체 적용 기준")
print("슬라이드 12 완료")

# ─────────────────────────────────────────────────────
# SLIDE 13 — 결론 & 다음 단계
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_shape(slide, Inches(0), Inches(0), W, H, fill_color=PRIMARY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.15), fill_color=ACCENT)
add_shape(slide, W - Inches(0.5), Inches(0), Inches(0.5), H, fill_color=ACCENT)
add_text(slide, "13", W - Inches(0.85), Inches(0.05), Inches(0.7), Inches(0.25),
         font_size=10, color=MUTED, align=PP_ALIGN.RIGHT)

add_text(slide, "결론 & 다음 단계", Inches(0.6), Inches(0.5), Inches(10), Inches(0.5),
         font_size=18, bold=True, color=ACCENT)

conclusions = [
    ("바이럴 루프가 핵심 성장 엔진",
     "K-factor 3.95, gift_received 유입 40% — 선물 경험이 선물을 부른다\n캠페인 예산의 70%를 시즌 D-7~D-3에 집중하면 최대 효율"),
    ("Champions·At Risk 2개 세그먼트가 최우선",
     "합산 GMV ₩16억 (전체 21%), 유저 수 5,193명(11.6%)\nVIP 보호 + 이탈 방지가 신규 고객 확보보다 ROAS 효율 높음"),
    ("30일 reciprocity 58% = 골든타임",
     "선물 수신 직후 30일이 발신 전환 핵심 구간\n수신 D+3~D+7 자동 리마인더 캠페인으로 Viral Loop 증폭 가능"),
]
for i, (title, body) in enumerate(conclusions):
    y = Inches(1.2 + i * 1.8)
    add_shape(slide, Inches(0.6), y, Inches(0.08), Inches(1.5),
              fill_color=[ACCENT, DANGER, SUCCESS][i])
    add_text(slide, title, Inches(0.85), y + Inches(0.05), Inches(10), Inches(0.45),
             font_size=14, bold=True, color=WHITE)
    add_text(slide, body, Inches(0.85), y + Inches(0.52), Inches(11.5), Inches(0.75),
             font_size=11, color=LIGHT)

# 다음 단계
add_text(slide, "다음 단계", Inches(0.6), Inches(6.6), Inches(3), Inches(0.3),
         font_size=11, bold=True, color=ACCENT)
next_steps = "1. STL period=30 재분석 [완료]  |  2. ROAS 감도분석 10/15/20원 [완료]  |  3. 분기별 K-factor 자동 모니터링  |  4. 시즌 D-7 캠페인 A/B 테스트 설계"
add_text(slide, next_steps, Inches(0.6), Inches(6.95), Inches(12.0), Inches(0.3),
         font_size=10, color=MUTED)

print("슬라이드 13 완료")

# ─────────────────────────────────────────────────────
# SLIDE 14 — Appendix: STL Decomposition (period=30)
# Reviewer 플래그 수정: period=7 → period=30 (월간 시즌성 포착)
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 14, "Appendix")
section_title(slide, "[수정] STL Decomposition period=30 — 월간 시즌성 성분 재산출")
takeaway_box(slide, "Reviewer 지적사항 반영: period=7(주간) → period=30(월간)으로 연간 시즌 이벤트 포착 강화")

# STL 설명 박스
stl_notes = [
    ("period=7 (이전)", "주간 주기만 포착. 빼빼로데이·설날 등 연간 이벤트 시즌성 성분 과소 추정"),
    ("period=30 (수정)", "월간 주기 포착. 시즌성 성분 표준편차 2,471,865원 → 3,208,555원 (+29.8%)"),
]
for i, (label, desc) in enumerate(stl_notes):
    y = Inches(1.45 + i * 0.65)
    col = MUTED if i == 0 else SUCCESS
    add_shape(slide, Inches(0.4), y, Inches(0.07), Inches(0.55), fill_color=col)
    add_text(slide, label, Inches(0.6), y + Inches(0.04), Inches(2.5), Inches(0.28),
             font_size=11, bold=True, color=col)
    add_text(slide, desc, Inches(3.2), y + Inches(0.04), Inches(9.8), Inches(0.5),
             font_size=10, color=PRIMARY)

ANALYSIS_CHART_DIR = os.path.join(BASE, "analysis", "charts")
add_picture_safe(slide, os.path.join(ANALYSIS_CHART_DIR, "layer1_stl_decomposition.png"),
                 Inches(0.3), Inches(2.85), Inches(12.7), Inches(4.2))
footer(slide, "주: STL = Seasonal-Trend decomposition using LOESS | robust=True (이상치 내성)")
print("슬라이드 14 완료 (Appendix: STL period=30)")

# ─────────────────────────────────────────────────────
# SLIDE 15 — Appendix: ROAS 감도분석
# Reviewer 플래그 수정: 15원 고정 → 10/15/20원 시나리오 비교
# ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
top_bar(slide, 15, "Appendix")
section_title(slide, "[수정] ROAS 감도분석 — 비용 10/15/20원 시나리오에서 결론 불변")
takeaway_box(slide, "Reviewer 지적사항 반영: 비용 가정이 2배 달라져도 상위 3개 세그먼트 ROAS 600x 이상 유지")

# 감도분석 수치 테이블
sens_header = ["세그먼트", "10원/건", "15원/건 (기본)", "20원/건", "결론"]
sens_rows = [
    ("Champions", "2,497x", "1,665x", "1,249x", "비용 무관 최고 ROAS"),
    ("At Risk", "1,213x", "809x", "606x", "안정적 고효율"),
    ("Need Attention", "1,205x", "803x", "602x", "예상수익 최대 (7,010만)"),
    ("Loyal", "1,142x", "761x", "571x", "절대 수익 1위 (9,940만)"),
]
col_widths = [3.2, 1.9, 2.5, 1.9, 3.0]
col_x = [0.5, 3.7, 5.6, 8.1, 10.0]

# 헤더
add_shape(slide, Inches(0.4), Inches(1.5), Inches(12.6), Inches(0.42), fill_color=PRIMARY)
for j, (cell, w, x) in enumerate(zip(sens_header, col_widths, col_x)):
    add_text(slide, cell, Inches(x), Inches(1.55), Inches(w - 0.1), Inches(0.32),
             font_size=10, bold=True, color=WHITE)

# 데이터 행
row_colors = [LIGHT, WHITE, LIGHT, WHITE]
for i, row in enumerate(sens_rows):
    y = Inches(1.95 + i * 0.58)
    add_shape(slide, Inches(0.4), y, Inches(12.6), Inches(0.52), fill_color=row_colors[i])
    for j, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
        # 15원 컬럼 강조
        bold = (j == 2)
        color = HIGHLIGHT if j == 2 else PRIMARY
        add_text(slide, cell, Inches(x), y + Inches(0.1), Inches(w - 0.1), Inches(0.35),
                 font_size=10, bold=bold, color=color)

# 감도분석 차트
add_picture_safe(slide, os.path.join(ANALYSIS_CHART_DIR, "layer5_roas_sensitivity.png"),
                 Inches(0.3), Inches(4.35), Inches(12.7), Inches(2.8))
footer(slide, "주: ROAS = (평균구매액 × CVR) / 비용 | 비용 범위 10~20원은 카카오 친구톡 실제 단가 범위")
print("슬라이드 15 완료 (Appendix: ROAS 감도분석)")

# ─────────────────────────────────────────────────────
# 저장
# ─────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"\nPPTX 저장 완료: {OUT_PATH}")
